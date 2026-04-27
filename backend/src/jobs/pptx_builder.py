from __future__ import annotations

from copy import deepcopy
import os
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches

from src.core.logging import get_logger
from src.core.settings import get_settings


logger = get_logger('appslides.backend.pptx')
_SETTINGS = get_settings()
_FONT_FALLBACK = (_SETTINGS.font_fallback or 'Cambria').strip()
_FONT_WHITELIST = {font.lower() for font in _SETTINGS.font_whitelist} if _SETTINGS.font_whitelist else set()
_FONTS_DIR = str(_SETTINGS.fonts_dir).strip()
_AVAILABLE_FONTS: set[str] = set()
_FONT_SUBSTITUTIONS = {
    'cambria': 'Caladea',
    'calibri': 'Carlito',
}


def _load_fonts_from_dir(fonts_dir: str) -> set[str]:
    if not fonts_dir:
        return set()
    names: list[str] = []
    try:
        for root, _dirs, files in os.walk(fonts_dir):
            for name in files:
                if not name.lower().endswith(('.ttf', '.otf', '.ttc')):
                    continue
                base = os.path.splitext(name)[0].replace('_', ' ').strip()
                if base:
                    names.append(base)
                    if '-' in base:
                        family = base.split('-', 1)[0].strip()
                        if family:
                            names.append(family)
    except OSError:
        return set()
    return {name.lower() for name in names if name}


if _FONTS_DIR:
    _AVAILABLE_FONTS = _load_fonts_from_dir(_FONTS_DIR)


def build_presentation(template_path: str, slides: list[dict[str, Any]], output_path: str) -> str:
    logger.info('PPTX build start: template=%s slides=%s', template_path, len(slides))
    presentation = Presentation(template_path)
    if len(presentation.slides) == 0:
        raise ValueError('Template has no slides')

    template_slides = list(presentation.slides)
    total_needed = len(slides)

    for index, slide_data in enumerate(slides):
        if index < len(template_slides):
            slide = template_slides[index]
        else:
            slide = _clone_slide(presentation, template_slides[-1])

        title_set, body_set, image_set = _apply_tokens(slide, slide_data)
        if not title_set:
            title = slide.shapes.title
            if title:
                _apply_text_to_shape(title, str(slide_data.get('title', '')), is_title=True)
        if not body_set:
            body = _find_body_placeholder(slide)
            if body:
                _apply_text_to_shape(body, str(slide_data.get('text', '')))
        if not image_set:
            image_path = slide_data.get('image_path')
            if image_path:
                _add_image(slide, str(image_path))

    if len(template_slides) > total_needed:
        for extra_slide in template_slides[total_needed:]:
            _remove_slide(presentation, extra_slide)

    presentation.save(output_path)
    logger.info('PPTX build done: output=%s', output_path)
    return output_path


def _find_body_placeholder(slide):
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.type in (
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.VERTICAL_BODY,
            PP_PLACEHOLDER.VERTICAL_OBJECT,
        ):
            return shape
    return None


def _add_image(slide, image_path: str) -> None:
    placeholder = _find_picture_placeholder(slide)
    if placeholder:
        try:
            placeholder.insert_picture(image_path)
            return
        except Exception:
            pass
    left = Inches(6.0)
    top = Inches(1.6)
    width = Inches(3.5)
    slide.shapes.add_picture(image_path, left, top, width=width)


def _find_picture_placeholder(slide):
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.type in (PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.OBJECT):
            return shape
    return None


def _remove_slide(presentation: Presentation, slide) -> None:
    slide_id_list = presentation.slides._sldIdLst  # pylint: disable=protected-access
    for item in list(slide_id_list):
        if int(item.get('id')) == slide.slide_id:
            presentation.part.drop_rel(item.rId)
            slide_id_list.remove(item)
            return


def _clone_slide(presentation: Presentation, source_slide):
    slide_layout = source_slide.slide_layout
    new_slide = presentation.slides.add_slide(slide_layout)
    for shape in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shape._element)
    for shape in source_slide.shapes:
        new_element = deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
    try:
        if source_slide._element.bg is not None and hasattr(new_slide._element, 'get_or_add_bg'):  # pylint: disable=protected-access
            new_slide._element.get_or_add_bg().append(deepcopy(source_slide._element.bg))  # pylint: disable=protected-access
    except Exception:
        pass
    return new_slide


def _apply_tokens(slide, slide_data: dict[str, Any]) -> tuple[bool, bool, bool]:
    title_set = False
    body_set = False
    image_set = False
    image_box = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text or ''
        if _has_token(text, 'IMAGE'):
            image_box = shape
            image_set = True
            continue
        if _has_token(text, 'TITLE'):
            _apply_text_to_shape(shape, str(slide_data.get('title', '')), is_title=True)
            title_set = True
        if _has_token(text, 'TEXT'):
            _apply_text_to_shape(shape, str(slide_data.get('text', '')))
            body_set = True

    if image_box and slide_data.get('image_path'):
        _replace_shape_with_image(slide, image_box, str(slide_data['image_path']))

    return title_set, body_set, image_set


def _has_token(text: str, token: str) -> bool:
    token = token.upper()
    markers = [f'##{token}', f'###{token}', f'{{{{{token}}}}}']
    return any(marker in text.upper() for marker in markers)


def _replace_shape_with_image(slide, shape, image_path: str) -> None:
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    shape._element.getparent().remove(shape._element)
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)


def _apply_text_to_shape(shape, text: str, is_title: bool = False) -> None:
    if not shape.has_text_frame:
        return
    spec = _get_font_spec(shape)
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.clear()
    lines = [line for line in text.splitlines()]
    if not lines:
        return
    text_frame.text = lines[0]
    for line in lines[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line

    max_size = 24
    if spec and spec.get('size') is not None:
        try:
            max_size = int(spec['size'].pt)
        except Exception:
            max_size = 24
    min_size = 14 if is_title else 10
    if is_title:
        max_size = min(max_size, 40)
    size = _calc_font_size(shape, text, max_size, min_size)
    if size <= min_size:
        text = _truncate_text_to_fit(shape, text, size)
        text_frame.clear()
        lines = [line for line in text.splitlines()]
        if lines:
            text_frame.text = lines[0]
            for line in lines[1:]:
                paragraph = text_frame.add_paragraph()
                paragraph.text = line

    _apply_font_size(text_frame, size)
    resolved_font = _resolve_font_name(spec.get('name') if spec else None)
    if resolved_font:
        _apply_font_to_text_frame(text_frame, {'name': resolved_font})

    if resolved_font:
        try:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            text_frame.fit_text(
                font_family=resolved_font,
                max_size=size,
                bold=bool(spec.get('bold')),
                italic=bool(spec.get('italic')),
            )
        except Exception:
            pass


def _get_font_spec(shape) -> dict[str, Any] | None:
    if not shape.has_text_frame:
        return None
    text_frame = shape.text_frame
    if not text_frame.paragraphs:
        return None
    paragraph = text_frame.paragraphs[0]
    if paragraph.runs:
        return _font_from_run(paragraph.runs[0])
    return _font_from_paragraph(paragraph)


def _font_from_run(run) -> dict[str, Any] | None:
    font = run.font
    if not font:
        return None
    return {
        'name': font.name,
        'size': font.size,
        'color': getattr(font.color, 'rgb', None) if font.color else None,
        'bold': font.bold,
        'italic': font.italic,
    }


def _font_from_paragraph(paragraph) -> dict[str, Any] | None:
    font = paragraph.font
    if not font:
        return None
    return {
        'name': font.name,
        'size': font.size,
        'color': getattr(font.color, 'rgb', None) if font.color else None,
        'bold': font.bold,
        'italic': font.italic,
    }


def _apply_font_to_text_frame(text_frame, spec: dict[str, Any] | None) -> None:
    if not spec:
        return
    for paragraph in text_frame.paragraphs:
        runs = paragraph.runs if paragraph.runs else _ensure_runs(paragraph)
        for run in runs:
            font = run.font
            if spec.get('name'):
                font.name = spec['name']
            if spec.get('size'):
                font.size = spec['size']
            if spec.get('color'):
                font.color.rgb = spec['color']
            if spec.get('bold') is not None:
                font.bold = spec['bold']
            if spec.get('italic') is not None:
                font.italic = spec['italic']


def _ensure_runs(paragraph):
    paragraph.text = paragraph.text
    return paragraph.runs


def _resolve_font_name(name: str | None) -> str | None:
    if not _FONT_FALLBACK and not name:
        return name
    if not name:
        return _FONT_FALLBACK
    if _AVAILABLE_FONTS:
        normalized = name.lower()
        if normalized in _AVAILABLE_FONTS:
            return name
        substitute = _FONT_SUBSTITUTIONS.get(normalized)
        if substitute and substitute.lower() in _AVAILABLE_FONTS:
            return substitute
        return _FONT_FALLBACK
    if _FONT_WHITELIST:
        return name if name.lower() in _FONT_WHITELIST else _FONT_FALLBACK
    return name


def _calc_font_size(shape, text: str, max_size: int, min_size: int) -> int:
    width_pt = shape.width / 12700
    height_pt = shape.height / 12700
    paragraphs = [paragraph.strip() for paragraph in text.splitlines() if paragraph.strip()]
    if not paragraphs:
        return max_size

    for size in range(max_size, min_size - 1, -1):
        line_height = size * 1.25
        total_lines = 0
        for paragraph in paragraphs:
            chars_per_line = max(1, int(width_pt / (0.65 * size)))
            total_lines += (len(paragraph) + chars_per_line - 1) // chars_per_line
        needed_height = total_lines * line_height
        if needed_height <= height_pt:
            return size
    return min_size


def _truncate_text_to_fit(shape, text: str, size: int) -> str:
    width_pt = shape.width / 12700
    height_pt = shape.height / 12700
    line_height = size * 1.25
    max_lines = max(1, int(height_pt / line_height))
    chars_per_line = max(1, int(width_pt / (0.65 * size)))
    max_chars = max_lines * chars_per_line
    flat = ' '.join(paragraph.strip() for paragraph in text.splitlines() if paragraph.strip())
    if len(flat) <= max_chars:
        return text
    trimmed = flat[: max(0, max_chars - 1)].rstrip()
    return trimmed + '…'


def _apply_font_size(text_frame, size: int) -> None:
    for paragraph in text_frame.paragraphs:
        runs = paragraph.runs if paragraph.runs else _ensure_runs(paragraph)
        for run in runs:
            run.font.size = size * 12700
