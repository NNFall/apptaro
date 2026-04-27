import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

TEMPLATES = [
    {
        'name': 'design_1.pptx',
        'bg': RGBColor(245, 248, 255),
        'title_color': RGBColor(20, 30, 60),
        'body_color': RGBColor(30, 30, 30),
        'title_font': 'Calibri',
        'body_font': 'Calibri',
    },
    {
        'name': 'design_2.pptx',
        'bg': RGBColor(255, 247, 235),
        'title_color': RGBColor(85, 35, 20),
        'body_color': RGBColor(50, 40, 30),
        'title_font': 'Arial',
        'body_font': 'Arial',
    },
    {
        'name': 'design_3.pptx',
        'bg': RGBColor(240, 255, 246),
        'title_color': RGBColor(10, 60, 40),
        'body_color': RGBColor(25, 45, 35),
        'title_font': 'Georgia',
        'body_font': 'Georgia',
    },
    {
        'name': 'design_4.pptx',
        'bg': RGBColor(248, 248, 248),
        'title_color': RGBColor(20, 20, 20),
        'body_color': RGBColor(50, 50, 50),
        'title_font': 'Verdana',
        'body_font': 'Verdana',
    },
]


def _apply_style(slide, title_text: str, body_text: str, spec: dict) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = spec['bg']

    title = slide.shapes.title
    if title:
        title.text = title_text
        for p in title.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = spec['title_font']
                r.font.size = Pt(36)
                r.font.color.rgb = spec['title_color']
                r.font.bold = True

    body = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == 2:
            body = shape
            break
    if body:
        body.text = body_text
        for p in body.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = spec['body_font']
                r.font.size = Pt(20)
                r.font.color.rgb = spec['body_color']


def main() -> None:
    out_dir = os.path.join('media', 'templates')
    os.makedirs(out_dir, exist_ok=True)
    for spec in TEMPLATES:
        prs = Presentation()
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        _apply_style(slide, 'Заголовок', 'Текст слайда', spec)
        out_path = os.path.join(out_dir, spec['name'])
        prs.save(out_path)
        print('created', out_path)


if __name__ == '__main__':
    main()
