from __future__ import annotations

from copy import deepcopy
import os
from typing import List, Dict, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches

from services.logger import get_logger
from config import load_config


logger = get_logger()
_CONFIG = load_config()
_FONT_FALLBACK = (_CONFIG.font_fallback or 'Cambria').strip()
_FONT_WHITELIST = {f.lower() for f in _CONFIG.font_whitelist} if _CONFIG.font_whitelist else set()
_FONTS_DIR = (_CONFIG.fonts_dir or '').strip()
_AVAILABLE_FONTS = set()


def _load_fonts_from_dir(fonts_dir: str) -> set[str]:
    if not fonts_dir:
        return set()
    names: list[str] = []
    try:
        for root, _dirs, files in os.walk(fonts_dir):
            for name in files:
                if not name.lower().endswith(('.ttf', '.otf', '.ttc')):
                    continue
                base = os.path.splitext(name)[0].replace('_', ' ').strip()
                if base:
                    names.append(base)
                    if '-' in base:
                        family = base.split('-', 1)[0].strip()
                        if family:
                            names.append(family)
    except OSError:
        return set()
    return {n.lower() for n in names if n}


if _FONTS_DIR:
    _AVAILABLE_FONTS = _load_fonts_from_dir(_FONTS_DIR)


def build_presentation(template_path: str, slides: List[Dict[str, str]], output_path: str) -> str:
    logger.info('PPTX build start: template=%s slides=%s', template_path, len(slides))
    prs = Presentation(template_path)
    if len(prs.slides) == 0:
        raise ValueError('Template has no slides')

    template_slides = list(prs.slides)
    total_needed = len(slides)

    for idx, slide_data in enumerate(slides):
        if idx < len(template_slides):
            slide = template_slides[idx]
        else:
            slide = _clone_slide(prs, template_slides[-1])

        title_set, body_set, image_set = _apply_tokens(slide, slide_data)

        if not title_set:
            title = slide.shapes.title
            if title:
                _apply_text_to_shape(title, slide_data.get('title', ''), is_title=True)

        if not body_set:
            body = _find_body_placeholder(slide)
            if body:
                _apply_text_to_shape(body, slide_data.get('text', ''))

        if not image_set:
            image_path = slide_data.get('image_path')
            if image_path:
                _add_image(slide, image_path)

    if len(template_slides) > total_needed:
        for extra_slide in template_slides[total_needed:]:
            _remove_slide(prs, extra_slide)

    prs.save(output_path)
    logger.info('PPTX build done: output=%s', output_path)
    return output_path


def _find_body_placeholder(slide):
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.type in (
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.VERTICAL_BODY,
            PP_PLACEHOLDER.VERTICAL_OBJECT,
        ):
            return shape
    return None


def _add_image(slide, image_path: str) -> None:
    placeholder = _find_picture_placeholder(slide)
    if placeholder:
        try:
            placeholder.insert_picture(image_path)
            return
        except Exception:  # noqa: BLE001
            pass
    left = Inches(6.0)
    top = Inches(1.6)
    width = Inches(3.5)
    slide.shapes.add_picture(image_path, left, top, width=width)


def _find_picture_placeholder(slide):
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.type in (PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.OBJECT):
            return shape
    return None


def _remove_slide(prs: Presentation, slide) -> None:
    slide_id_list = prs.slides._sldIdLst  # pylint: disable=protected-access
    for sld in list(slide_id_list):
        if int(sld.get('id')) == slide.slide_id:
            prs.part.drop_rel(sld.rId)
            slide_id_list.remove(sld)
            return


def _clone_slide(prs: Presentation, source_slide):
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shape._element)
    for shape in source_slide.shapes:
        new_el = deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    try:
        if source_slide._element.bg is not None:  # pylint: disable=protected-access
            if hasattr(new_slide._element, 'get_or_add_bg'):
                new_slide._element.get_or_add_bg().append(  # pylint: disable=protected-access
                    deepcopy(source_slide._element.bg)
                )
    except Exception:  # noqa: BLE001
        # Безопасно пропускаем фон, если python-pptx не поддерживает get_or_add_bg
        pass
    return new_slide


def _apply_tokens(slide, slide_data: Dict[str, str]) -> Tuple[bool, bool, bool]:
    title_set = False
    body_set = False
    image_set = False
    image_box = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text or ''
        if _has_token(text, 'IMAGE'):
            image_box = shape
            image_set = True
            continue
        if _has_token(text, 'TITLE'):
            _apply_text_to_shape(shape, slide_data.get('title', ''), is_title=True)
            title_set = True
        if _has_token(text, 'TEXT'):
            _apply_text_to_shape(shape, slide_data.get('text', ''))
            body_set = True

    if image_box and slide_data.get('image_path'):
        _replace_shape_with_image(slide, image_box, slide_data['image_path'])

    return title_set, body_set, image_set


def _has_token(text: str, token: str) -> bool:
    token = token.upper()
    markers = [f'##{token}', f'###{token}', f'{{{{{token}}}}}']
    return any(m in text.upper() for m in markers)


def _replace_shape_with_image(slide, shape, image_path: str) -> None:
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    shape._element.getparent().remove(shape._element)
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)


def _apply_text_to_shape(shape, text: str, is_title: bool = False) -> None:
    if not shape.has_text_frame:
        return
    spec = _get_font_spec(shape)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.clear()
    lines = [line for line in text.splitlines()]
    if not lines:
        return
    tf.text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line

    max_size = 24
    if spec and spec.get('size') is not None:
        try:
            max_size = int(spec['size'].pt)
        except Exception:  # noqa: BLE001
            max_size = 24
    if is_title:
        max_size = min(max_size, 40)
        min_size = 14
    else:
        min_size = 10
    size = _calc_font_size(shape, text, max_size, min_size)
    if size <= min_size:
        text = _truncate_text_to_fit(shape, text, size)
        tf.clear()
        lines = [line for line in text.splitlines()]
        if lines:
            tf.text = lines[0]
            for line in lines[1:]:
                p = tf.add_paragraph()
                p.text = line
    _apply_font_size(tf, size)
    resolved_font = _resolve_font_name(spec.get('name') if spec else None)
    if resolved_font:
        _apply_font_to_text_frame(tf, {'name': resolved_font})

    if resolved_font:
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.fit_text(
                font_family=resolved_font,
                max_size=size,
                bold=bool(spec.get('bold')),
                italic=bool(spec.get('italic')),
            )
        except Exception:  # noqa: BLE001
            pass


def _extract_template_style(prs: Presentation) -> dict:
    style = {}
    if len(prs.slides) == 0:
        return style
    slide = prs.slides[0]
    bg_color = _get_background_color(slide)
    if bg_color:
        style['background'] = bg_color
    title = slide.shapes.title
    if title:
        style['title_font'] = _get_font_spec(title)
    body = _find_body_placeholder(slide)
    if body:
        style['body_font'] = _get_font_spec(body)
    return style


def _get_background_color(slide) -> RGBColor | None:
    fill = slide.background.fill
    if not fill or fill.type != MSO_FILL.SOLID:
        return None
    color = fill.fore_color
    if color.type == MSO_COLOR_TYPE.RGB:
        return color.rgb
    return None


def _apply_background(slide, color: RGBColor | None) -> None:
    if not color:
        return
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _get_font_spec(shape) -> dict | None:
    if not shape.has_text_frame:
        return None
    tf = shape.text_frame
    if not tf.paragraphs:
        return None
    p = tf.paragraphs[0]
    if p.runs:
        r = p.runs[0]
        return _font_from_run(r)
    return _font_from_paragraph(p)


def _font_from_run(run) -> dict | None:
    font = run.font
    if not font:
        return None
    return {
        'name': font.name,
        'size': font.size,
        'color': getattr(font.color, 'rgb', None) if font.color else None,
        'bold': font.bold,
        'italic': font.italic,
    }


def _font_from_paragraph(p) -> dict | None:
    font = p.font
    if not font:
        return None
    return {
        'name': font.name,
        'size': font.size,
        'color': getattr(font.color, 'rgb', None) if font.color else None,
        'bold': font.bold,
        'italic': font.italic,
    }


def _apply_font_to_shape(shape, spec: dict | None) -> None:
    if not spec or not shape.has_text_frame:
        return
    _apply_font_to_text_frame(shape.text_frame, spec)


def _apply_font_to_text_frame(tf, spec: dict | None) -> None:
    if not spec:
        return
    for p in tf.paragraphs:
        if p.runs:
            runs = p.runs
        else:
            p.text = p.text
            runs = p.runs
        for run in runs:
            font = run.font
            if spec.get('name'):
                font.name = spec['name']
            if spec.get('size'):
                font.size = spec['size']
            if spec.get('color'):
                font.color.rgb = spec['color']
            if spec.get('bold') is not None:
                font.bold = spec['bold']
            if spec.get('italic') is not None:
                font.italic = spec['italic']


def _resolve_font_name(name: str | None) -> str | None:
    if not _FONT_FALLBACK and not name:
        return name
    if not name:
        return _FONT_FALLBACK
    if _AVAILABLE_FONTS:
        return name if name.lower() in _AVAILABLE_FONTS else _FONT_FALLBACK
    if _FONT_WHITELIST:
        return name if name.lower() in _FONT_WHITELIST else _FONT_FALLBACK
    return name


def _calc_font_size(shape, text: str, max_size: int, min_size: int) -> int:
    width_pt = shape.width / 12700
    height_pt = shape.height / 12700

    paragraphs = [p.strip() for p in text.splitlines() if p.strip()]
    if not paragraphs:
        return max_size

    for size in range(max_size, min_size - 1, -1):
        line_height = size * 1.25
        total_lines = 0
        for p in paragraphs:
            chars_per_line = max(1, int(width_pt / (0.65 * size)))
            total_lines += (len(p) + chars_per_line - 1) // chars_per_line
        needed_height = total_lines * line_height
        if needed_height <= height_pt:
            return size
    return min_size


def _truncate_text_to_fit(shape, text: str, size: int) -> str:
    width_pt = shape.width / 12700
    height_pt = shape.height / 12700
    line_height = size * 1.25
    max_lines = max(1, int(height_pt / line_height))
    chars_per_line = max(1, int(width_pt / (0.65 * size)))
    max_chars = max_lines * chars_per_line
    flat = ' '.join(p.strip() for p in text.splitlines() if p.strip())
    if len(flat) <= max_chars:
        return text
    trimmed = flat[: max(0, max_chars - 1)].rstrip()
    return trimmed + '…'


def _apply_font_size(tf, size: int) -> None:
    for p in tf.paragraphs:
        if p.runs:
            runs = p.runs
        else:
            p.text = p.text
            runs = p.runs
        for run in runs:
            run.font.size = size * 12700
